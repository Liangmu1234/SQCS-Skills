#!/usr/bin/env python3
"""Build a hybrid editable PPTX from layout.json using python-pptx."""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def parse_hex(value: str | None, default: str = "#000000") -> RGBColor:
    value = (value or default).strip()
    if not value.startswith("#"):
        value = "#" + value
    try:
        return RGBColor(int(value[1:3], 16), int(value[3:5], 16), int(value[5:7], 16))
    except Exception:
        d = default.lstrip("#")
        return RGBColor(int(d[0:2], 16), int(d[2:4], 16), int(d[4:6], 16))


def inch(v: float):
    return Inches(float(v))


def apply_line(shape, el):
    line = el.get("line")
    if line and str(line).lower() != "none":
        shape.line.color.rgb = parse_hex(line, "#D0D5DD")
        shape.line.width = Pt(float(el.get("line_width", 1) or 1))
    else:
        shape.line.fill.background()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("layout", type=Path)
    parser.add_argument("--out", type=Path, default=Path("reconstructed.pptx"))
    args = parser.parse_args()

    layout = json.loads(args.layout.read_text(encoding="utf-8"))
    slide_cfg = layout.get("slide", {})
    width = float(slide_cfg.get("width", 13.333))
    height = float(slide_cfg.get("height", 7.5))
    theme = layout.get("theme", {})
    default_font = theme.get("font", "Microsoft YaHei")

    prs = Presentation()
    prs.slide_width = inch(width)
    prs.slide_height = inch(height)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = parse_hex(theme.get("background", "#FFFFFF"), "#FFFFFF")

    base_dir = args.layout.parent

    for el in layout.get("elements", []):
        etype = el.get("type")
        name = el.get("name")
        if etype == "rect":
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if float(el.get("radius", 0) or 0) > 0 else MSO_SHAPE.RECTANGLE
            shp = slide.shapes.add_shape(shape_type, inch(el.get("x", 0)), inch(el.get("y", 0)), inch(el.get("w", 0)), inch(el.get("h", 0)))
            if name:
                shp.name = str(name)
            fill = el.get("fill")
            if fill and str(fill).lower() != "none":
                shp.fill.solid()
                shp.fill.fore_color.rgb = parse_hex(fill, "#FFFFFF")
                # python-pptx transparency support is limited across versions.
            else:
                shp.fill.background()
            apply_line(shp, el)

        elif etype == "text":
            tb = slide.shapes.add_textbox(inch(el.get("x", 0)), inch(el.get("y", 0)), inch(el.get("w", 0)), inch(el.get("h", 0)))
            if name:
                tb.name = str(name)
            tf = tb.text_frame
            tf.clear()
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "center": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }.get(str(el.get("valign", "top")).lower(), MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            p.text = str(el.get("text", ""))
            p.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }.get(str(el.get("align", "left")).lower(), PP_ALIGN.LEFT)
            for run in p.runs:
                run.font.name = el.get("font", default_font)
                run.font.size = Pt(float(el.get("font_size", 12) or 12))
                run.font.bold = bool(el.get("bold"))
                run.font.italic = bool(el.get("italic"))
                run.font.color.rgb = parse_hex(el.get("color", "#000000"), "#000000")

        elif etype == "line":
            shp = slide.shapes.add_connector(1, inch(el.get("x1", 0)), inch(el.get("y1", 0)), inch(el.get("x2", 0)), inch(el.get("y2", 0)))
            if name:
                shp.name = str(name)
            shp.line.color.rgb = parse_hex(el.get("color", "#000000"), "#000000")
            shp.line.width = Pt(float(el.get("width", 1) or 1))

        elif etype == "image":
            img_path = Path(str(el.get("path", "")))
            if not img_path.is_absolute():
                img_path = base_dir / img_path
            if img_path.exists():
                shp = slide.shapes.add_picture(str(img_path), inch(el.get("x", 0)), inch(el.get("y", 0)), inch(el.get("w", 0)), inch(el.get("h", 0)))
                if name:
                    shp.name = str(name)
            else:
                # Fallback placeholder rectangle if asset is missing.
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(el.get("x", 0)), inch(el.get("y", 0)), inch(el.get("w", 0)), inch(el.get("h", 0)))
                shp.name = f"missing_asset_{name or 'image'}"
                shp.fill.solid()
                shp.fill.fore_color.rgb = parse_hex("#FEE4E2")
                shp.line.color.rgb = parse_hex("#D92D20")
                shp.text = f"Missing asset: {el.get('path','')}"

        elif etype == "table":
            rows = int(el.get("rows", len(el.get("data", [])) or 1))
            cols = int(el.get("cols", max((len(r) for r in el.get("data", [])), default=1)))
            shp = slide.shapes.add_table(rows, cols, inch(el.get("x", 0)), inch(el.get("y", 0)), inch(el.get("w", 0)), inch(el.get("h", 0)))
            if name:
                shp.name = str(name)
            tbl = shp.table
            data = el.get("data", [])
            for r in range(rows):
                for c in range(cols):
                    cell = tbl.cell(r, c)
                    cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = default_font
                            run.font.size = Pt(float(el.get("font_size", 10) or 10))
                            run.font.color.rgb = parse_hex(el.get("text_color", "#101828"), "#101828")

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
